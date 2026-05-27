"""
L0: PPT 빌드 후 슬라이드 레이아웃 overflow 자동 검증 — 경계 이탈 즉시 차단
PostToolUse hook on Bash — build_*_ppt.py 실행 시만 발동
"""
import sys, json, os, glob, time, re

try:
    raw = sys.stdin.read() or '{}'
    data = json.loads(raw)
    cmd = data.get('tool_input', {}).get('command', '')

    # build_*_ppt.py 포함한 명령만 처리
    if 'build_' not in cmd or '_ppt.py' not in cmd:
        sys.exit(0)

    # 스크립트 경로 추출
    m = re.search(r'["\']?([A-Za-z]:[/\\][^"\';\n]+build_\w+_ppt\.py)["\']?', cmd)
    if not m:
        m = re.search(r'([^\s"\']+build_\w+_ppt\.py)', cmd)
    script_path = m.group(1).strip('"\'') if m else None

    if not script_path:
        sys.exit(0)

    base_dir = os.path.dirname(os.path.abspath(script_path.replace('/', os.sep).replace('\\', os.sep)))
    if not os.path.isdir(base_dir):
        sys.exit(0)

    # 최근 120초 내 수정된 .pptx 찾기
    now = time.time()
    pptx_files = glob.glob(os.path.join(base_dir, '*.pptx'))
    recent = sorted(
        [f for f in pptx_files if now - os.path.getmtime(f) < 120],
        key=lambda f: os.path.getmtime(f), reverse=True
    )
    if not recent:
        sys.exit(0)

    from pptx import Presentation
    from pptx.util import Inches

    SW  = Inches(13.33)   # 슬라이드 폭
    SH  = Inches(7.5)     # 슬라이드 높이
    TOL = Inches(0.01)    # 허용 오차

    errors = []
    target = recent[0]  # 가장 최근 파일만 검증
    try:
        prs = Presentation(target)
        fname = os.path.basename(target)
        for si, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                try:
                    r = shape.left + shape.width
                    b = shape.top + shape.height
                    if r > SW + TOL:
                        errors.append(f"슬라이드{si} '{shape.name}': 우측이탈 r={r/914400:.3f}\" (한계:{SW/914400:.2f}\")")
                    if b > SH + TOL:
                        errors.append(f"슬라이드{si} '{shape.name}': 하단이탈 b={b/914400:.3f}\" (한계:{SH/914400:.2f}\")")
                except Exception:
                    pass
    except Exception as e:
        print(json.dumps({"hookSpecificOutput": {"hookEventName": "PostToolUse",
            "additionalContext": f"PPTX 검증 실패: {e}"}}))
        sys.exit(0)

    if errors:
        lines = errors[:10]
        extra = f"\n... 외 {len(errors)-10}개" if len(errors) > 10 else ""
        msg = f"⚠️ PPTX OVERFLOW {len(errors)}건!\n" + "\n".join(lines) + extra
        print(json.dumps({
            "systemMessage": msg,
            "hookSpecificOutput": {"hookEventName": "PostToolUse", "additionalContext": msg}
        }))
        sys.exit(2)
    else:
        msg = f"✓ PPTX 레이아웃 검증 통과 — {os.path.basename(target)}"
        print(json.dumps({"hookSpecificOutput": {"hookEventName": "PostToolUse", "additionalContext": msg}}))
        sys.exit(0)

except Exception:
    sys.exit(0)  # 훅 오류는 작업 차단하지 않음
